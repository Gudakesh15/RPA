#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build GermanVocabBot_Presentation.pptx
Run: python3 make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x3C, 0x8C)
NAVY2   = RGBColor(0x28, 0x4B, 0x9E)   # lighter navy for secondary boxes
ORANGE  = RGBColor(0xE8, 0x63, 0x0A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0xF2, 0xF4, 0xF8)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
MID     = RGBColor(0x55, 0x55, 0x72)
LBLUE   = RGBColor(0xD6, 0xE4, 0xF7)
GREEN   = RGBColor(0x1E, 0x8A, 0x44)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GOLD    = RGBColor(0xF5, 0xA6, 0x23)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0)):
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Pt
    s = slide.shapes.add_shape(1, x, y, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def tb(slide, text, x, y, w, h,
       size=14, color=DARK, bold=False, italic=False,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(size)
    run.font.color.rgb = color
    run.font.bold     = bold
    run.font.italic   = italic
    return box


def header(slide, title, subtitle=None):
    """Navy header bar with white title + optional subtitle."""
    rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    rect(slide, 0, Inches(1.25), W, Inches(0.055), fill=ORANGE)
    tb(slide, title,
       Inches(0.5), Inches(0.1), Inches(11.5), Inches(0.75),
       size=34, color=WHITE, bold=True)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.5), Inches(0.82), Inches(11.5), Inches(0.38),
           size=14, color=LBLUE, italic=True)


def card(slide, x, y, w, h, title, body,
         accent=NAVY, bg=LGREY, t_size=15, b_size=12):
    rect(slide, x, y, w, h, fill=bg)
    rect(slide, x, y, w, Inches(0.07), fill=accent)
    tb(slide, title,
       x+Inches(0.17), y+Inches(0.12), w-Inches(0.34), Inches(0.42),
       size=t_size, color=NAVY, bold=True)
    tb(slide, body,
       x+Inches(0.17), y+Inches(0.57), w-Inches(0.34), h-Inches(0.7),
       size=b_size, color=MID)


# ── Build slides ─────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

rect(s, 0, 0,               W, Inches(0.14), fill=ORANGE)
rect(s, 0, H-Inches(0.14),  W, Inches(0.14), fill=ORANGE)

# Subtle background stripe
rect(s, 0, Inches(2.8), W, Inches(2.5), fill=RGBColor(0x14, 0x2E, 0x72))

# Main title
tb(s, "German Vocabulary Bot",
   Inches(0.8), Inches(1.5), Inches(11.73), Inches(1.5),
   size=62, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Orange divider
rect(s, Inches(4.2), Inches(3.2), Inches(4.93), Inches(0.055), fill=ORANGE)

# Tagline
tb(s, "Automating daily language learning with UiPath · Python · NotebookLM",
   Inches(0.8), Inches(3.38), Inches(11.73), Inches(0.65),
   size=20, color=LBLUE, align=PP_ALIGN.CENTER)

# Course badge
rect(s, Inches(4.5), Inches(4.3), Inches(4.33), Inches(0.55), fill=NAVY2)
tb(s, "RPA Seminar — Masters Semester 2",
   Inches(4.5), Inches(4.32), Inches(4.33), Inches(0.5),
   size=13, color=GOLD, align=PP_ALIGN.CENTER)

# Team
tb(s, "Rodrigo  ·  Avi  ·  [Name]",
   Inches(0.8), Inches(5.15), Inches(11.73), Inches(0.5),
   size=19, color=WHITE, align=PP_ALIGN.CENTER)

tb(s, "github.com/Gudakesh15/RPA",
   Inches(0.8), Inches(6.8), Inches(11.73), Inches(0.38),
   size=12, color=RGBColor(0x60, 0x80, 0xA8), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Our Story
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "Our Story", "Three non-native German speakers with the same daily struggle")

# Large quote box
rect(s, Inches(0.45), Inches(1.55), Inches(12.43), Inches(1.65), fill=LBLUE)
rect(s, Inches(0.45), Inches(1.55), Inches(0.1), Inches(1.65), fill=NAVY)
tb(s, "“We search German words all day long…",
   Inches(0.7), Inches(1.68), Inches(11.8), Inches(0.65),
   size=26, color=NAVY, bold=True, italic=True)
tb(s, "  and forget every single one by evening.”",
   Inches(0.7), Inches(2.25), Inches(11.8), Inches(0.65),
   size=26, color=NAVY, bold=True, italic=True)

# Three persona cards
personas = [
    ("\U0001f393  The Student",
     "Attending German lectures and seminars, constantly looking up academic "
     "vocabulary — Bedeutung, zuständig, Antrag — "
     "and losing them all in browser history."),
    ("\U0001f3d9️  The City Dweller",
     "Reading menus, official letters, street signs and texts in German every "
     "day. New words appear everywhere — none of them stick."),
    ("\U0001f4ac  The Daily Learner",
     "Watching German TV, chatting with colleagues, searching words "
     "mid-conversation. The searches pile up. The review never happens."),
]
cw = Inches(3.97)
ch = Inches(2.45)
for i, (title, body) in enumerate(personas):
    x = Inches(0.45) + i * (cw + Inches(0.27))
    card(s, x, Inches(3.4), cw, ch, title, body,
         accent=ORANGE if i == 1 else NAVY, bg=LGREY, t_size=16, b_size=12)

tb(s, "The problem: vocabulary is searched constantly, but there is no automatic way to capture and review it.",
   Inches(0.45), Inches(6.1), Inches(12.43), Inches(0.45),
   size=14, color=MID, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "The Problem", "Vocabulary is searched every day — but never revisited")

# Left stat panel
rect(s, Inches(0.4), Inches(1.55), Inches(5.6), Inches(5.5), fill=LGREY)
tb(s, "5–20", Inches(0.4), Inches(1.85), Inches(5.6), Inches(2.0),
   size=100, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, "German words searched per day",
   Inches(0.5), Inches(3.85), Inches(5.4), Inches(0.6),
   size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
tb(s, "across Google, Duden, DeepL, Linguee, DWDS",
   Inches(0.5), Inches(4.45), Inches(5.4), Inches(0.4),
   size=12, color=MID, align=PP_ALIGN.CENTER)

rect(s, Inches(0.7), Inches(4.95), Inches(5.0), Inches(0.055), fill=ORANGE)

tb(s, "0  reviewed",
   Inches(0.4), Inches(5.1), Inches(5.6), Inches(0.75),
   size=46, color=RED, bold=True, align=PP_ALIGN.CENTER)
tb(s, "at the end of the day",
   Inches(0.5), Inches(5.85), Inches(5.4), Inches(0.4),
   size=14, color=MID, align=PP_ALIGN.CENTER)

# Right: four pain points
pain = [
    ("\U0001f4f1  Multi-device problem",
     "Searches happen on phone, tablet and laptop.\nNo single place collects them all."),
    ("⏱️  No time to review",
     "By the time you sit down to study, you can’t\nremember what you searched for."),
    ("\U0001f4cb  Manual effort always fails",
     "Copy-pasting words into flashcard apps sounds\ngood — but never actually happens consistently."),
    ("\U0001f501  Repetition without progress",
     "The same words get searched again and again.\nLearning never sticks because review never happens."),
]
for i, (title, body) in enumerate(pain):
    y = Inches(1.55) + i * Inches(1.38)
    rect(s, Inches(6.25), y, Inches(6.65), Inches(1.2), fill=LGREY)
    rect(s, Inches(6.25), y, Inches(0.07), Inches(1.2), fill=RED)
    tb(s, title, Inches(6.45), y+Inches(0.1), Inches(6.2), Inches(0.38),
       size=14, color=RED, bold=True)
    tb(s, body, Inches(6.45), y+Inches(0.5), Inches(6.2), Inches(0.65),
       size=12, color=MID)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Solution
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "The Solution", "A fully automated pipeline: from browser history to AI revision video")

steps = [
    ("\U0001f310", "Google\nHistory",     "All searches\nacross every device\nvia Google sync"),
    ("\U0001f916", "UiPath\nOrchestrator","Scrapes history,\ntriggers scripts,\nnavigates apps"),
    ("\U0001f40d", "Python\nFilter",      "Detects and isolates\nGerman vocabulary\nwords"),
    ("\U0001f4c4", "Source\nFile",        "Clean plain-text\nvocabulary list\nfor NotebookLM"),
    ("\U0001f3ac", "AI Video\nOverview",  "NotebookLM creates\nyour personalised\nrevision episode"),
]

BW = Inches(2.12)
BH = Inches(3.35)
AW = Inches(0.34)
total = len(steps)*BW + (len(steps)-1)*AW
sx = (W - total) / 2

for i, (icon, title, desc) in enumerate(steps):
    x = sx + i*(BW+AW)
    y = Inches(1.9)

    is_last = (i == len(steps)-1)
    bg   = NAVY  if is_last else (LBLUE if i % 2 == 0 else LGREY)
    tc   = WHITE if is_last else NAVY
    dc   = RGBColor(0xB0,0xC4,0xDE) if is_last else MID

    rect(s, x, y, BW, BH, fill=bg)

    tb(s, icon, x, y+Inches(0.22), BW, Inches(0.7),
       size=34, align=PP_ALIGN.CENTER)
    tb(s, title, x, y+Inches(1.0), BW, Inches(0.72),
       size=16, color=tc, bold=True, align=PP_ALIGN.CENTER)
    tb(s, desc, x+Inches(0.12), y+Inches(1.78), BW-Inches(0.24), Inches(1.35),
       size=11, color=dc, align=PP_ALIGN.CENTER)

    if not is_last:
        ax = x + BW
        tb(s, "→", ax, y+Inches(1.2), AW, Inches(0.6),
           size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

tb(s, "One UiPath workflow orchestrates the entire flow — from scraping to video generation.",
   Inches(0.5), Inches(5.6), Inches(12.33), Inches(0.45),
   size=14, color=MID, italic=True, align=PP_ALIGN.CENTER)

# Google sync callout
rect(s, Inches(0.35), Inches(6.18), Inches(4.5), Inches(0.85), fill=LBLUE)
tb(s, "\U0001f4a1  Key insight: Google syncs history across phone + laptop.\n"
      "    UiPath captures everything you searched on every device.",
   Inches(0.5), Inches(6.22), Inches(4.2), Inches(0.75),
   size=11, color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — UiPath in Action
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "UiPath in Action", "22 activities — one browser session — fully automated end to end")

left = [
    ("①  Scrape",    "Opens myactivity.google.com\nExtracts full search history table"),
    ("②  Store",     "Writes to Excel via\nUse Excel File activity"),
    ("③  Process",   "Triggers Python script\nvia Start Process activity"),
    ("④  Navigate",  "Goes to notebooklm.google.com\nin the same Chrome window"),
]
right = [
    ("⑤  Create",    "Clicks ‘New Notebook’\nautomatically"),
    ("⑥  Upload",    "Opens Windows file picker\nUploads vocabulary source file"),
    ("⑦  Customise", "Clicks Video Overview\nTypes personalised AI prompt"),
    ("⑧  Complete",  "Clicks Generate\nShows ‘Process completed’ dialog"),
]

RH = Inches(1.3)
RW = Inches(5.85)
for i, (title, body) in enumerate(left):
    y = Inches(1.55) + i*(RH+Inches(0.12))
    rect(s, Inches(0.35), y, RW, RH, fill=LGREY)
    rect(s, Inches(0.35), y, Inches(0.07), RH, fill=NAVY)
    tb(s, title, Inches(0.55), y+Inches(0.1), RW-Inches(0.25), Inches(0.42),
       size=16, color=NAVY, bold=True)
    tb(s, body, Inches(0.55), y+Inches(0.55), RW-Inches(0.25), Inches(0.65),
       size=12, color=MID)

for i, (title, body) in enumerate(right):
    y = Inches(1.55) + i*(RH+Inches(0.12))
    rect(s, Inches(7.13), y, RW, RH, fill=LGREY)
    rect(s, Inches(7.13), y, Inches(0.07), RH, fill=ORANGE)
    tb(s, title, Inches(7.33), y+Inches(0.1), RW-Inches(0.25), Inches(0.42),
       size=16, color=NAVY, bold=True)
    tb(s, body, Inches(7.33), y+Inches(0.55), RW-Inches(0.25), Inches(0.65),
       size=12, color=MID)

# Centre connector
rect(s, Inches(6.34), Inches(1.55), Inches(0.65), Inches(5.85), fill=LGREY)
tb(s, "→", Inches(6.32), Inches(4.1), Inches(0.68), Inches(0.55),
   size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, "Google\nActivity", Inches(6.3), Inches(1.7), Inches(0.72), Inches(0.9),
   size=9, color=MID, align=PP_ALIGN.CENTER)
tb(s, "Notebook\nLM", Inches(6.3), Inches(5.9), Inches(0.72), Inches(0.9),
   size=9, color=MID, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Challenges & Solutions
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "Challenges We Faced", "Real debugging experience building the UiPath workflow")

challenges = [
    ("Python not found by UiPath",
     "Start Process with 'python' failed — Python not in system PATH when UiPath launches it.",
     "Full exe path: C:\\Python313\\python.exe\n+ Chr(34) quoting for paths that contain spaces"),
    ("Dynamic UI element IDs",
     "NotebookLM uses Angular Material — element IDs like 'mat-input-1' change on every page load, breaking selectors.",
     "Positional selector: <webctrl tag='TEXTAREA' idx='3' />\nStable regardless of Angular’s ID assignment"),
    ("Wrong data format assumption",
     "Expected URL columns from myactivity.google.com — UiPath Extract Table Data returned a single text column instead.",
     "Rewrote Script 1 to parse the format:\n‘Searched for [term]\\n[time] • Details’ directly"),
    ("UI in a foreign language",
     "NotebookLM displayed ‘Generar’ (Spanish) instead of ‘Generate’ — strict text selector failed completely.",
     "UiPath Fuzzy Selector targeting aaname=' Generar '\nWorks regardless of NotebookLM’s display language"),
]

CW = Inches(5.95)
CH = Inches(2.5)
for i, (problem, explain, fix) in enumerate(challenges):
    row, col = divmod(i, 2)
    x = Inches(0.35) + col*(CW+Inches(0.38))
    y = Inches(1.55) + row*(CH+Inches(0.18))

    rect(s, x, y, CW, CH, fill=LGREY)
    rect(s, x, y, CW, Inches(0.07), fill=RED)

    tb(s, "⚠️  " + problem,
       x+Inches(0.17), y+Inches(0.12), CW-Inches(0.34), Inches(0.42),
       size=14, color=RED, bold=True)
    tb(s, explain,
       x+Inches(0.17), y+Inches(0.57), CW-Inches(0.34), Inches(0.62),
       size=11, color=MID)

    rect(s, x+Inches(0.17), y+Inches(1.25), CW-Inches(0.34), Inches(0.04), fill=GREEN)

    tb(s, "✓  " + fix,
       x+Inches(0.17), y+Inches(1.34), CW-Inches(0.34), Inches(1.05),
       size=11, color=GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MVP Complete
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "MVP — What We Built", "End-to-end pipeline running and recorded")

checklist = [
    "Scrape Google My Activity (cross-device, phone + laptop)",
    "Write full search history to Excel via UiPath",
    "Trigger Python script automatically from UiPath",
    "Detect German words: ä / ö / ü / ß character detection",
    "Generate clean source file for NotebookLM",
    "Auto-create a new NotebookLM notebook",
    "Upload vocabulary file via Windows file picker",
    "Trigger AI Video Overview in NotebookLM",
    "Type personalised learning prompt automatically",
    "Generate AI video — confirm with ‘Process completed’",
]
for i, item in enumerate(checklist):
    y = Inches(1.6) + i*Inches(0.5)
    rect(s, Inches(0.4), y+Inches(0.07), Inches(0.3), Inches(0.3), fill=GREEN)
    tb(s, "✓", Inches(0.4), y+Inches(0.04), Inches(0.3), Inches(0.35),
       size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, item, Inches(0.82), y, Inches(5.7), Inches(0.45),
       size=13, color=DARK)

# Right panel — outputs
rect(s, Inches(7.15), Inches(1.55), Inches(5.75), Inches(5.55), fill=NAVY)
tb(s, "Files Produced Each Run",
   Inches(7.35), Inches(1.68), Inches(5.35), Inches(0.5),
   size=17, color=WHITE, bold=True)

outputs = [
    ("raw_searches.xlsx",       "Full Google search history"),
    ("vocabulary_words.xlsx",   "Filtered German word list"),
    ("notebooklm_source.txt",   "Upload file for NotebookLM"),
    ("NotebookLM Notebook",     "New notebook created automatically"),
    ("AI Video Overview",       "Personalised revision video \U0001f3ac"),
    ("digest_YYYYMMDD.pdf",     "PDF backup (secondary path)"),
]
for i, (fname, desc) in enumerate(outputs):
    y = Inches(2.35) + i*Inches(0.72)
    rect(s, Inches(7.35), y, Inches(5.15), Inches(0.65), fill=NAVY2)
    tb(s, fname, Inches(7.5), y+Inches(0.04), Inches(4.85), Inches(0.3),
       size=13, color=GOLD, bold=True)
    tb(s, desc,  Inches(7.5), y+Inches(0.34), Inches(4.85), Inches(0.25),
       size=11, color=LBLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What's Next
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "What’s Next", "Making UiPath the core of every step — more automation, smarter filtering")

future = [
    ("\U0001f501", "UiPath Data\nFiltering Loop",
     "UiPath reads vocabulary Excel, loops through every row, applies its own filter rules "
     "and builds the upload file — less Python, more UiPath intelligence"),
    ("\U0001f4d6", "DWDS\nEnrichment",
     "UiPath opens dwds.de for each word, extracts the definition and an example sentence, "
     "writes it back to Excel — richer content for the video"),
    ("⚡", "Conditional\nLogic",
     "If zero German words found, UiPath skips NotebookLM and shows a smart notification "
     "instead of uploading an empty file"),
    ("\U0001f4c5", "Daily\nScheduling",
     "Windows Task Scheduler runs the bot every morning — your revision video is ready "
     "before the first lecture of the day"),
    ("\U0001f4ca", "Run Log\n& Audit Trail",
     "UiPath appends each run to a log Excel: date, words found, notebook name, status "
     "— track vocabulary growth over weeks"),
    ("\U0001f4e7", "Email\nDigest",
     "UiPath sends a weekly summary email with word count, the full vocabulary list, "
     "and a direct link to the NotebookLM notebook"),
]

CW = Inches(3.93)
CH = Inches(2.1)
for i, (icon, title, body) in enumerate(future):
    row, col = divmod(i, 3)
    x = Inches(0.35) + col*(CW+Inches(0.235))
    y = Inches(1.6) + row*(CH+Inches(0.22))
    accent = NAVY if row == 0 else ORANGE
    rect(s, x, y, CW, CH, fill=LGREY)
    rect(s, x, y, CW, Inches(0.07), fill=accent)
    tb(s, icon + "  " + title,
       x+Inches(0.17), y+Inches(0.14), CW-Inches(0.34), Inches(0.55),
       size=15, color=NAVY, bold=True)
    tb(s, body,
       x+Inches(0.17), y+Inches(0.72), CW-Inches(0.34), CH-Inches(0.85),
       size=11, color=MID)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Full Vision
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
header(s, "The Full Vision", "A daily autonomous language learning companion")

# Centre hub
rect(s, Inches(4.55), Inches(2.65), Inches(4.23), Inches(2.0), fill=NAVY)
rect(s, Inches(4.55), Inches(2.65), Inches(4.23), Inches(0.08), fill=ORANGE)
tb(s, "Daily Automated\nVocabulary Digest",
   Inches(4.55), Inches(2.75), Inches(4.23), Inches(1.8),
   size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Satellite cards
sats = [
    (Inches(0.3),  Inches(1.65), "\U0001f4f1  Multi-device", "Phone + laptop searches\ncaptured via Google sync"),
    (Inches(0.3),  Inches(3.35), "\U0001f310  Multi-language", "French, Spanish, Italian\n— same architecture"),
    (Inches(0.3),  Inches(5.05), "\U0001f9e0  LLM Enrichment", "Definitions + examples\nadded via API before upload"),
    (Inches(9.8),  Inches(1.65), "\U0001f4dd  Anki Flashcards", "Auto-generated deck\nfrom vocabulary list"),
    (Inches(9.8),  Inches(3.35), "\U0001f4ca  Orchestrator", "Track vocabulary growth\nover weeks and months"),
    (Inches(9.8),  Inches(5.05), "\U0001f4e7  Email Digest", "Weekly summary email\nevery Sunday morning"),
]

for x, y, title, body in sats:
    side = "left" if x < Inches(5) else "right"
    bg = LBLUE if side == "left" else LGREY
    rect(s, x, y, Inches(3.15), Inches(1.55), fill=bg)
    tb(s, title, x+Inches(0.17), y+Inches(0.1), Inches(2.8), Inches(0.45),
       size=14, color=NAVY, bold=True)
    tb(s, body, x+Inches(0.17), y+Inches(0.6), Inches(2.8), Inches(0.85),
       size=12, color=MID)
    # Arrow
    if side == "left":
        tb(s, "→", x+Inches(3.18), y+Inches(0.45), Inches(0.45), Inches(0.55),
           size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    else:
        tb(s, "←", x-Inches(0.48), y+Inches(0.45), Inches(0.45), Inches(0.55),
           size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
rect(s, 0, 0,              W, Inches(0.14), fill=ORANGE)
rect(s, 0, H-Inches(0.14), W, Inches(0.14), fill=ORANGE)
rect(s, 0, Inches(2.9), W, Inches(2.2), fill=RGBColor(0x14, 0x2E, 0x72))

tb(s, "Thank You",
   Inches(0.8), Inches(1.1), Inches(11.73), Inches(1.55),
   size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(4.5), Inches(3.1), Inches(4.33), Inches(0.06), fill=ORANGE)

tb(s, "Questions?",
   Inches(0.8), Inches(3.25), Inches(11.73), Inches(0.75),
   size=30, color=LBLUE, align=PP_ALIGN.CENTER)

summary = [
    "\U0001f310  Any German word you search — captured automatically, across every device",
    "\U0001f916  UiPath orchestrates the entire pipeline without any manual steps",
    "\U0001f3ac  NotebookLM turns your daily vocabulary into a personalised revision video",
]
for i, line in enumerate(summary):
    tb(s, line,
       Inches(1.8), Inches(4.3) + i*Inches(0.72), Inches(9.73), Inches(0.6),
       size=17, color=WHITE, align=PP_ALIGN.CENTER)

tb(s, "github.com/Gudakesh15/RPA",
   Inches(0.8), Inches(6.85), Inches(11.73), Inches(0.38),
   size=13, color=RGBColor(0x60, 0x80, 0xA8), align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
OUT = "/Users/aviking/Desktop/Masters/Sem 2/RPA/GermanVocabBot_Presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
